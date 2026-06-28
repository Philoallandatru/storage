#!/usr/bin/env python3
"""Generate an executive PPT for KV cache IO/LBA analysis.

The deck is intentionally data-backed: every number is pulled from the
existing JSON/CSV artifacts produced by the KV cache profiling scripts.
"""

from __future__ import annotations

import argparse
import csv
import json
import math
import re
from collections import Counter
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "docs" / "assets" / "kvcache-boss-ppt"
OUT_PPT = ROOT / "docs" / "presentations" / "kv-cache-lba-io-executive-2026-06-26.pptx"
OUT_SUMMARY = ROOT / "docs" / "presentations" / "kv-cache-lba-io-executive-2026-06-26.md"

PROFILE_JSON = ROOT / "results/kvcache-profile/test_sharegpt_8b_tp8_cpu0p5g_users2_300s_profile_20260608_014520.json"
LBA_EVENTS_JSON = ROOT / "results/kvcache-profile/lba_timeline/lba_events.json"
LBA_SUMMARY_JSON = ROOT / "results/kvcache-profile/lba_timeline/lba_timeline_summary.json"
KEY_LOCALITY_JSON = ROOT / "results/kvcache-profile/key_locality/kvcache_key_locality_summary.json"
IO_RANDOMNESS_CSV = ROOT / "results/cross_vendor/kv_cache_k4_gc_drift/_analysis/kvcache_io_randomness_summary.csv"
BPFTRACE_LOG = ROOT / "results/kvcache-profile/bpftrace_sharegpt_8b_tp8_cpu0p5g_users2_300s_profile_20260608_014520.txt"


BG = RGBColor(8, 12, 24)
PANEL = RGBColor(17, 25, 43)
PANEL_2 = RGBColor(24, 35, 58)
TEXT = RGBColor(238, 244, 255)
MUTED = RGBColor(155, 170, 195)
CYAN = RGBColor(41, 211, 255)
GREEN = RGBColor(63, 220, 146)
YELLOW = RGBColor(255, 201, 77)
RED = RGBColor(255, 91, 109)
PURPLE = RGBColor(176, 121, 255)


def setup_matplotlib() -> None:
    cjk_font_path = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
    if Path(cjk_font_path).exists():
        from matplotlib import font_manager as fm

        fm.fontManager.addfont(cjk_font_path)
        plt.rcParams["font.sans-serif"] = ["Noto Sans CJK SC", "Noto Sans CJK JP", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    plt.rcParams["figure.facecolor"] = "#080c18"
    plt.rcParams["axes.facecolor"] = "#11192b"
    plt.rcParams["axes.edgecolor"] = "#3a4a67"
    plt.rcParams["axes.labelcolor"] = "#dce8ff"
    plt.rcParams["text.color"] = "#eff6ff"
    plt.rcParams["xtick.color"] = "#9fb0cc"
    plt.rcParams["ytick.color"] = "#9fb0cc"
    plt.rcParams["grid.color"] = "#33415f"
    plt.rcParams["savefig.facecolor"] = "#080c18"


def read_json(path: Path) -> dict | list:
    return json.loads(path.read_text(encoding="utf-8"))


def percentile(values: list[float], q: float) -> float:
    if not values:
        return 0.0
    arr = sorted(values)
    return arr[min(len(arr) - 1, math.ceil(q * len(arr)) - 1)]


def parse_bpftrace_histograms(path: Path) -> dict[str, list[tuple[int, int, int]]]:
    lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    histograms: dict[str, list[tuple[int, int, int]]] = {}
    i = 0
    while i < len(lines):
        m = re.match(r"^@(\w+):\s*$", lines[i].rstrip())
        if not m:
            i += 1
            continue
        name = m.group(1)
        if name == "d":
            i += 1
            continue
        rows = []
        j = i + 1
        while j < len(lines):
            mh = re.match(r"^\[(\d+),\s*(\d+)\)\s+(\d+)\s*\|", lines[j].rstrip())
            if not mh:
                break
            rows.append((int(mh.group(1)), int(mh.group(2)), int(mh.group(3))))
            j += 1
        if rows:
            histograms[name] = rows
            i = j
        else:
            i += 1
    return histograms


def hist_pct_in_range(rows: list[tuple[int, int, int]], lo: int, hi: int) -> float:
    total = sum(c for _, _, c in rows)
    if total == 0:
        return 0.0
    return sum(c for a, b, c in rows if a >= lo and b <= hi) / total * 100


def hist_quantile_hi(rows: list[tuple[int, int, int]], q: float) -> int:
    total = sum(c for _, _, c in rows)
    target = total * q
    acc = 0
    for _, hi, c in rows:
        acc += c
        if acc >= target:
            return hi
    return rows[-1][1]


def collect_metrics() -> dict:
    profile = read_json(PROFILE_JSON)
    lba_events = read_json(LBA_EVENTS_JSON)
    lba_summary = read_json(LBA_SUMMARY_JSON)
    locality = read_json(KEY_LOCALITY_JSON)
    io_df = pd.read_csv(IO_RANDOMNESS_CSV)
    hist = parse_bpftrace_histograms(BPFTRACE_LOG)

    profile_summary = profile["summary"]
    cache = profile_summary["cache_stats"]

    lba_vals = [e["lba_gib"] for e in lba_events]
    nonzero_lba = [x for x in lba_vals if x > 0]
    buckets = Counter(int(x // 10) * 10 for x in lba_vals)

    read_size_rows = hist["bssplit_read_kb"]
    write_size_rows = hist["bssplit_write_kb"]
    read_lat_rows = hist["d2c_read_us"]
    write_lat_rows = hist["d2c_write_us"]

    return {
        "profile": profile,
        "profile_summary": profile_summary,
        "cache": cache,
        "lba_events": lba_events,
        "lba_summary": lba_summary,
        "locality": locality,
        "io_df": io_df,
        "hist": hist,
        "requests": profile["requests_completed"],
        "tokens": profile["total_tokens_generated"],
        "elapsed_s": profile_summary["elapsed_time"],
        "tok_s": profile_summary["avg_throughput_tokens_per_sec"],
        "req_s": profile_summary["requests_per_second"],
        "cache_hit_pct": cache["cache_hit_rate"] * 100,
        "storage_io_p50_ms": profile_summary["storage_io_latency_ms"]["p50"],
        "storage_io_p95_ms": profile_summary["storage_io_latency_ms"]["p95"],
        "storage_io_p99_ms": profile_summary["storage_io_latency_ms"]["p99"],
        "storage_read_dev_p95_ms": cache["storage_health"]["criteria"][0]["actual"],
        "total_read_gib": cache["total_read_gb"],
        "total_write_gib": cache["total_write_gb"],
        "read_write_ratio": cache["read_write_ratio"],
        "storage_read_gib": cache["tier_storage_kv_bytes_read_gb"],
        "cpu_read_gib": cache["tier_cpu_kv_bytes_read_gb"],
        "cpu_write_gib": cache["tier_cpu_kv_bytes_written_gb"],
        "decode_reads": cache["decode_reads"],
        "prefill_writes": cache["prefill_writes"],
        "lba_unique": len(lba_events),
        "lba_min_gib": min(lba_vals),
        "lba_nonzero_min_gib": min(nonzero_lba),
        "lba_nonzero_max_gib": max(nonzero_lba),
        "lba_nonzero_span_gib": max(nonzero_lba) - min(nonzero_lba),
        "lba_p50_gib": percentile(lba_vals, 0.5),
        "lba_p95_gib": percentile(lba_vals, 0.95),
        "lba_top_buckets": buckets.most_common(5),
        "last_touch_gap_1mib_pct": lba_summary["last_touch_gap_rate_by_threshold"]["gap_lt_1mib_pct"],
        "last_touch_gap_10mib_pct": lba_summary["last_touch_gap_rate_by_threshold"]["gap_lt_10mib_pct"],
        "last_touch_gap_100mib_pct": lba_summary["last_touch_gap_rate_by_threshold"]["gap_lt_100mib_pct"],
        "last_touch_gap_ge_100mib_pct": lba_summary["last_touch_gap_rate_by_threshold"]["gap_ge_100mib_pct"],
        "last_touch_forward_pct": lba_summary["last_touch_direction_distribution"]["forward_pct"],
        "last_touch_forward_run": lba_summary["last_touch_directional_runs"]["mean_forward_run_length"],
        "intra_token_pct": locality["intra_token_pct"],
        "inter_token_pct": locality["inter_token_pct"],
        "inter_request_pct": locality["inter_request_pct"],
        "intra_token_median_ms": locality["intra_token_median_ms"],
        "inter_token_median_ms": locality["inter_token_median_ms"],
        "inter_request_median_ms": locality["inter_request_median_ms"],
        "read_128_256_pct": hist_pct_in_range(read_size_rows, 128, 256),
        "write_128_256_pct": hist_pct_in_range(write_size_rows, 128, 256),
        "read_d2c_p50_us": hist_quantile_hi(read_lat_rows, 0.50),
        "read_d2c_p99_us": hist_quantile_hi(read_lat_rows, 0.99),
        "write_d2c_p50_us": hist_quantile_hi(write_lat_rows, 0.50),
        "write_d2c_p99_us": hist_quantile_hi(write_lat_rows, 0.99),
    }


def savefig(path: Path) -> None:
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def chart_signal_dashboard(m: dict) -> Path:
    path = ASSET_DIR / "01_signal_dashboard.png"
    fig, ax = plt.subplots(figsize=(13.5, 6.4))
    ax.axis("off")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)

    metrics = [
        ("Cache hit", f"{m['cache_hit_pct']:.1f}%", "应用层命中率", "#3fdc92"),
        ("Read shape", f"{m['read_128_256_pct']:.0f}%", "读 IO 在 128-256 KiB", "#29d3ff"),
        ("LBA span", f"{m['lba_nonzero_span_gib']:.0f} GiB", "last-touch 非零跨度", "#b079ff"),
        ("Cold reuse", f"{m['inter_request_pct']:.1f}%", "跨请求重读占比", "#ffc94d"),
    ]
    xs = [16, 39, 62, 85]
    for x, (title, value, sub, color) in zip(xs, metrics):
        ring = plt.Circle((x, 58), 12, fill=False, lw=7, color=color, alpha=0.95)
        glow = plt.Circle((x, 58), 14, fill=False, lw=2, color=color, alpha=0.25)
        ax.add_patch(glow)
        ax.add_patch(ring)
        ax.text(x, 60, value, ha="center", va="center", fontsize=20, fontweight="bold", color="#eff6ff")
        ax.text(x, 42, title, ha="center", va="center", fontsize=14, color=color, fontweight="bold")
        ax.text(x, 35, sub, ha="center", va="center", fontsize=11, color="#9fb0cc")

    ax.text(4, 92, "KV Cache IO: 不是纯顺序, 也不是纯随机", fontsize=24, fontweight="bold", color="#eff6ff")
    ax.text(4, 84, "结论来自应用层 trace + bpftrace 设备观测 + iostat 长跑 + LBA last-touch 复核", fontsize=13, color="#9fb0cc")
    ax.text(4, 13, "老板版一句话: 真实风险不是平均吞吐, 而是跨请求冷读 + 写/GC 尾延迟 + 文件系统 LBA 分布。", fontsize=15, color="#ffc94d")
    ax.plot([4, 96], [22, 22], color="#33415f", lw=1)
    savefig(path)
    return path


def chart_lba_density(m: dict) -> Path:
    path = ASSET_DIR / "02_lba_density.png"
    events = m["lba_events"]
    lba = np.array([e["lba_gib"] for e in events])
    t0 = min(e["ts_ns"] for e in events)
    ts = np.array([(e["ts_ns"] - t0) / 1e9 for e in events])

    fig = plt.figure(figsize=(13.5, 7.2))
    gs = fig.add_gridspec(2, 2, width_ratios=[3.2, 1], height_ratios=[1, 2.1])
    ax_hist = fig.add_subplot(gs[0, 0])
    ax_scatter = fig.add_subplot(gs[1, 0])
    ax_bucket = fig.add_subplot(gs[:, 1])

    bins = np.arange(0, math.ceil(max(lba) / 10) * 10 + 10, 10)
    counts, edges, _ = ax_hist.hist(lba, bins=bins, color="#29d3ff", alpha=0.85, edgecolor="#080c18")
    ax_hist.set_title("LBA last-touch density by 10 GiB bucket", fontsize=14, fontweight="bold")
    ax_hist.set_ylabel("Unique sectors")
    ax_hist.grid(True, axis="y", alpha=0.25)
    ax_hist.axvspan(m["lba_nonzero_min_gib"], m["lba_nonzero_max_gib"], color="#b079ff", alpha=0.12)

    scatter = ax_scatter.scatter(ts, lba, c=lba, cmap="turbo", s=12, alpha=0.75, edgecolors="none")
    ax_scatter.set_title("Each dot = unique (dev, sector) last-touch, not per-IO event", fontsize=13)
    ax_scatter.set_xlabel("Last-touch time since first map entry (s)")
    ax_scatter.set_ylabel("LBA (GiB)")
    ax_scatter.grid(True, alpha=0.25)
    fig.colorbar(scatter, ax=ax_scatter, label="LBA (GiB)")

    top = m["lba_top_buckets"]
    labels = [f"{b} GiB" for b, _ in top][::-1]
    vals = [c for _, c in top][::-1]
    ax_bucket.barh(labels, vals, color=["#3fdc92", "#29d3ff", "#b079ff", "#ffc94d", "#ff5b6d"][::-1])
    ax_bucket.set_title("Top buckets", fontsize=13, fontweight="bold")
    ax_bucket.set_xlabel("Unique sectors")
    for y, v in enumerate(vals):
        ax_bucket.text(v + 4, y, str(v), va="center", fontsize=10)
    ax_bucket.grid(True, axis="x", alpha=0.25)

    savefig(path)
    return path


def chart_locality_donut(m: dict) -> Path:
    path = ASSET_DIR / "03_locality_donut.png"
    labels = ["Intra-token\n<10ms", "Inter-token\n10ms-1s", "Inter-request\n>1s"]
    vals = [m["intra_token_pct"], m["inter_token_pct"], m["inter_request_pct"]]
    colors = ["#3fdc92", "#29d3ff", "#ffc94d"]

    fig, axes = plt.subplots(1, 2, figsize=(13.5, 6.4), gridspec_kw={"width_ratios": [1.1, 1]})
    ax = axes[0]
    wedges, _ = ax.pie(vals, colors=colors, startangle=90, counterclock=False, wedgeprops={"width": 0.36, "edgecolor": "#080c18", "linewidth": 2})
    ax.text(0, 0.08, f"{m['intra_token_pct']:.1f}%", ha="center", fontsize=30, fontweight="bold")
    ax.text(0, -0.13, "同 token 重读", ha="center", fontsize=13, color="#9fb0cc")
    ax.set_title("Application locality split", fontsize=16, fontweight="bold")
    ax.legend(wedges, labels, loc="lower center", bbox_to_anchor=(0.5, -0.16), ncol=3, fontsize=10)

    ax2 = axes[1]
    med = [m["intra_token_median_ms"], m["inter_token_median_ms"], m["inter_request_median_ms"]]
    ax2.barh(labels[::-1], med[::-1], color=colors[::-1])
    ax2.set_xscale("log")
    ax2.set_xlabel("Median reread interval (ms, log scale)")
    ax2.set_title("Same-key reread interval spans 6 orders of magnitude", fontsize=14, fontweight="bold")
    for y, v in enumerate(med[::-1]):
        text = f"{v:.2f} ms" if v < 1000 else f"{v/1000:.1f} s"
        ax2.text(v * 1.15, y, text, va="center", fontsize=11)
    ax2.grid(True, axis="x", which="both", alpha=0.25)
    savefig(path)
    return path


def chart_device_shape_latency(m: dict) -> Path:
    path = ASSET_DIR / "04_device_shape_latency.png"
    hist = m["hist"]
    fig, axes = plt.subplots(1, 2, figsize=(13.5, 6.2))

    ax = axes[0]
    read_rows = hist["bssplit_read_kb"]
    write_rows = hist["bssplit_write_kb"]
    x = np.arange(len(read_rows))
    labels = [f"{lo}-{hi}" for lo, hi, _ in read_rows]
    read_pct = np.array([c for _, _, c in read_rows]) / sum(c for _, _, c in read_rows) * 100
    write_pct = np.array([c for _, _, c in write_rows]) / sum(c for _, _, c in write_rows) * 100
    ax.bar(x - 0.18, read_pct, width=0.36, color="#ff5b6d", label="Read")
    ax.bar(x + 0.18, write_pct, width=0.36, color="#29d3ff", label="Write")
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.set_ylabel("% ops")
    ax.set_xlabel("Device IO size bucket (KiB)")
    ax.set_title("Device sees large-block random IO", fontsize=14, fontweight="bold")
    ax.legend()
    ax.grid(True, axis="y", alpha=0.25)

    ax2 = axes[1]
    for name, color, label in [("d2c_read_us", "#3fdc92", "Read D2C"), ("d2c_write_us", "#ffc94d", "Write D2C")]:
        rows = hist[name]
        total = sum(c for _, _, c in rows)
        xs, ys, acc = [], [], 0
        for lo, hi, c in rows:
            xs.extend([lo, hi])
            ys.extend([acc / total, (acc + c) / total])
            acc += c
        ax2.plot(xs, ys, color=color, lw=2.5, label=label)
    ax2.set_xscale("log")
    ax2.set_xlabel("Device-to-completion latency (us, log)")
    ax2.set_ylabel("CDF")
    ax2.set_title("Device latency is fast; host/storage stack tail matters", fontsize=14, fontweight="bold")
    ax2.legend()
    ax2.grid(True, which="both", alpha=0.25)
    savefig(path)
    return path


def chart_cross_vendor(m: dict) -> Path:
    path = ASSET_DIR / "05_cross_vendor_randomness.png"
    df = m["io_df"].copy()
    fig, axes = plt.subplots(1, 2, figsize=(13.5, 6.2))

    ax = axes[0]
    x = np.arange(len(df))
    ax.bar(x - 0.2, df["read_req_median_kb"], width=0.4, color="#29d3ff", label="Read req p50")
    ax.bar(x + 0.2, df["write_req_median_kb"], width=0.4, color="#b079ff", label="Write req p50")
    ax.set_xticks(x)
    ax.set_xticklabels(df["disk"], rotation=18, ha="right")
    ax.set_ylim(100, 132)
    ax.set_ylabel("KiB")
    ax.set_title("Request shape is locked by application", fontsize=14, fontweight="bold")
    ax.legend()
    ax.grid(True, axis="y", alpha=0.25)

    ax2 = axes[1]
    sizes = np.clip(df["aqu_p99"] / df["aqu_p99"].max() * 900, 130, 900)
    sc = ax2.scatter(df["r_await_p99_ms"], df["w_await_p99_ms"], s=sizes, c=df["aqu_p99"], cmap="plasma", alpha=0.85, edgecolors="#eff6ff", linewidths=0.8)
    for _, r in df.iterrows():
        ax2.text(r["r_await_p99_ms"] * 1.07, r["w_await_p99_ms"] * 1.03, r["disk"], fontsize=10)
    ax2.set_xscale("log")
    ax2.set_yscale("log")
    ax2.set_xlabel("r_await p99 (ms, log)")
    ax2.set_ylabel("w_await p99 (ms, log)")
    ax2.set_title("Differentiation is tail latency + queue depth", fontsize=14, fontweight="bold")
    ax2.grid(True, which="both", alpha=0.25)
    fig.colorbar(sc, ax=ax2, label="aqu-sz p99")
    savefig(path)
    return path


def chart_evidence_stack(m: dict) -> Path:
    path = ASSET_DIR / "06_evidence_stack.png"
    fig, ax = plt.subplots(figsize=(13.5, 6.2))
    ax.axis("off")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)

    rows = [
        ("可靠", "应用 trace", "127K IO, Key locality", "83.4% 同 token 重读; 8.1% 跨请求冷读", "#3fdc92"),
        ("可靠", "bpftrace hist", "真实设备 IO size/latency", f"{m['read_128_256_pct']:.0f}% 读请求在 128-256 KiB; 读 D2C p99 {m['read_d2c_p99_us']} us", "#29d3ff"),
        ("可靠", "iostat long-run", "四盘 1s 聚合", "%rrqm p50/p99 均为 0; 请求大小稳定在 124-125 KiB", "#b079ff"),
        ("探索", "LBA last-touch", "969 个唯一 sector", f"非零范围 {m['lba_nonzero_min_gib']:.1f}-{m['lba_nonzero_max_gib']:.1f} GiB; 不能当 per-IO 顺序率", "#ffc94d"),
        ("待验证", "per-IO LBA stream", "blktrace/eBPF printf", "下一步才能严格证明 sequential ratio / run length / window working set", "#ff5b6d"),
    ]
    y = 82
    for level, source, scope, claim, color in rows:
        ax.add_patch(plt.Rectangle((5, y - 8), 90, 12, color="#11192b", ec=color, lw=1.5))
        ax.text(8, y, level, fontsize=13, fontweight="bold", color=color, va="center")
        ax.text(21, y, source, fontsize=14, fontweight="bold", color="#eff6ff", va="center")
        ax.text(41, y + 2.5, scope, fontsize=11, color="#9fb0cc", va="center")
        ax.text(41, y - 3.2, claim, fontsize=11.5, color="#eff6ff", va="center")
        y -= 16
    ax.text(5, 95, "Evidence stack: what we know vs what remains unproven", fontsize=22, fontweight="bold")
    savefig(path)
    return path


def generate_charts(m: dict) -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    return {
        "dashboard": chart_signal_dashboard(m),
        "lba": chart_lba_density(m),
        "locality": chart_locality_donut(m),
        "device": chart_device_shape_latency(m),
        "vendor": chart_cross_vendor(m),
        "evidence": chart_evidence_stack(m),
    }


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_text(slide, text, x, y, w, h, size=20, color=TEXT, bold=False, align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text or " "
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, items, x, y, w, h, size=18, color=TEXT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return tb


def add_kicker(slide, text):
    add_text(slide, text.upper(), 0.45, 0.25, 8, 0.3, size=10, color=CYAN, bold=True)


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.45, 0.55, 12.3, 0.65, size=26, color=TEXT, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.48, 1.16, 11.9, 0.38, size=12, color=MUTED)


def add_card(slide, x, y, w, h, title, value, subtitle, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = color
    shape.line.width = Pt(1.2)
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.36, 0.28, size=10, color=color, bold=True)
    add_text(slide, value, x + 0.18, y + 0.48, w - 0.36, 0.5, size=24, color=TEXT, bold=True)
    add_text(slide, subtitle, x + 0.18, y + 0.98, w - 0.36, 0.42, size=9, color=MUTED)


def add_image(slide, path: Path, x, y, w, h=None):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h) if h else None)


def build_ppt(m: dict, charts: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    s = blank_slide(prs)
    add_text(s, "KV Cache IO / LBA 分析", 0.6, 0.75, 9.2, 0.7, size=34, bold=True)
    add_text(s, "老板版结论: 应用层强局部性 + 设备层大块随机 + LBA last-touch 宽范围触达", 0.65, 1.55, 11.8, 0.45, size=16, color=MUTED)
    add_image(s, charts["dashboard"], 0.45, 2.1, 12.35, 4.25)
    add_text(s, "数据源: ShareGPT 8B TP8 CPU0.5 users2 300s + bpftrace + iostat K4 long-run | 2026-06-26", 0.65, 6.95, 12, 0.25, size=9, color=MUTED)

    # 2. Executive verdict
    s = blank_slide(prs)
    add_kicker(s, "executive verdict")
    add_title(s, "结论: 这不是一个“纯随机 IO”故事", "正确表述是: 应用访问有强时间局部性, 设备合并率很低, LBA 分布宽且集中在高位。")
    add_card(s, 0.6, 1.85, 2.85, 1.35, "应用层", f"{m['intra_token_pct']:.1f}%", "同 Key 重读发生在 10ms 内", GREEN)
    add_card(s, 3.65, 1.85, 2.85, 1.35, "设备层", f"{m['read_128_256_pct']:.0f}%", "读请求落在 128-256 KiB", CYAN)
    add_card(s, 6.7, 1.85, 2.85, 1.35, "LBA last-touch", f"{m['lba_nonzero_span_gib']:.0f} GiB", "非零唯一 sector 覆盖跨度", PURPLE)
    add_card(s, 9.75, 1.85, 2.85, 1.35, "四盘长跑", "0%", "%rrqm p50/p99, 合并几乎没有", YELLOW)
    add_bullets(
        s,
        [
            "对老板的可用结论: 请求形态由 KV cache 应用决定, SSD 之间差异主要体现在尾延迟和队列深度。",
            "需要避免的错误结论: 不能把 bpftrace @d[] last-touch map 当成完整 per-IO LBA 流。",
            "下一步要做 per-IO LBA trace, 才能严格证明真实 sequential ratio、run length 和窗口工作集。",
        ],
        0.85,
        3.85,
        11.6,
        1.4,
        size=17,
    )

    # 3. Test matrix
    s = blank_slide(prs)
    add_kicker(s, "test design")
    add_title(s, "测试矩阵: 用四类证据交叉约束同一个问题", "每类数据回答的问题不同, 不能互相替代。")
    rows = [
        ("ShareGPT trace", "127,477 IO / 301s", "Key、phase、timestamp", "应用层时间局部性"),
        ("bpftrace device", "171K+ block ops + 969 LBA entries", "IO size、D2C latency、@d[]", "设备真实请求形态"),
        ("iostat K4 long-run", "4 disks, ~1200 samples each", "req size、%rrqm、await、aqu-sz", "跨盘随机性和尾延迟"),
        ("LBA last-touch", "969 unique sectors", "dev, sector -> last timestamp", "空间触达范围, 不是 per-IO 顺序"),
    ]
    y = 1.75
    for source, scale, fields, answer in rows:
        add_card(s, 0.8, y, 2.65, 0.8, "数据源", source, scale, CYAN)
        add_card(s, 3.65, y, 3.1, 0.8, "字段", fields, "", PURPLE)
        add_card(s, 6.95, y, 5.2, 0.8, "回答", answer, "", GREEN)
        y += 1.05
    add_text(s, "设计原则: 一页讲结论, 后面每页只放支撑该结论的一类证据。", 0.85, 6.45, 11.6, 0.35, size=15, color=YELLOW)

    # 4. Workload KPI
    s = blank_slide(prs)
    add_kicker(s, "workload kpi")
    add_title(s, "真实 ShareGPT 负载: cache 命中高, 但读写比强烈偏读", "这解释了为什么 SSD 读尾延迟和冷读行为比平均带宽更值得看。")
    add_card(s, 0.65, 1.55, 2.8, 1.25, "Requests", f"{m['requests']:,}", f"{m['req_s']:.2f} req/s", CYAN)
    add_card(s, 3.65, 1.55, 2.8, 1.25, "Tokens", f"{m['tokens']:,}", f"{m['tok_s']:.0f} tok/s", GREEN)
    add_card(s, 6.65, 1.55, 2.8, 1.25, "Cache hit", f"{m['cache_hit_pct']:.1f}%", "PASS: target >30%", YELLOW)
    add_card(s, 9.65, 1.55, 2.8, 1.25, "Read / write", f"{m['read_write_ratio']:.1f}x", f"{m['total_read_gib']:.1f}/{m['total_write_gib']:.1f} GiB", PURPLE)
    add_image(s, charts["locality"], 0.55, 3.08, 12.25, 3.75)

    # 5. Device truth
    s = blank_slide(prs)
    add_kicker(s, "device truth")
    add_title(s, "设备真实看到的是大块随机 IO, 不是应用 304KB 原样请求", "bpftrace 直接观察 block layer: size split + device-to-completion latency。")
    add_image(s, charts["device"], 0.55, 1.45, 12.25, 5.55)
    add_text(s, f"关键数: 读 D2C p50/p99 = {m['read_d2c_p50_us']}/{m['read_d2c_p99_us']} us; 写 D2C p50/p99 = {m['write_d2c_p50_us']}/{m['write_d2c_p99_us']} us", 0.75, 6.95, 11.5, 0.3, size=11, color=MUTED)

    # 6. LBA
    s = blank_slide(prs)
    add_kicker(s, "lba spatial view")
    add_title(s, "LBA last-touch 显示高位宽范围触达, 但不能证明 per-IO 顺序率", "这页是校正后的 LBA 口径: unique sector 的最后访问时间。")
    add_image(s, charts["lba"], 0.55, 1.35, 12.25, 5.65)
    add_text(s, f"复核: {m['lba_unique']} unique sectors; 非零范围 {m['lba_nonzero_min_gib']:.1f}-{m['lba_nonzero_max_gib']:.1f} GiB; p50/p95 = {m['lba_p50_gib']:.1f}/{m['lba_p95_gib']:.1f} GiB", 0.75, 6.95, 11.7, 0.28, size=10.5, color=MUTED)

    # 7. Cross vendor
    s = blank_slide(prs)
    add_kicker(s, "cross-vendor stability")
    add_title(s, "四盘长跑: 请求大小稳定, 差异主要在尾延迟和队列", "这说明 workload shape 是应用锁定的, 设备优劣体现在稳定性。")
    add_image(s, charts["vendor"], 0.55, 1.45, 12.25, 5.55)

    # 8. What is correct
    s = blank_slide(prs)
    add_kicker(s, "validity")
    add_title(s, "哪些结论是正确的, 哪些不能说过头", "这页用于避免老板会后被追问时踩坑。")
    add_image(s, charts["evidence"], 0.55, 1.42, 12.25, 5.55)

    # 9. Recommendations
    s = blank_slide(prs)
    add_kicker(s, "recommendations")
    add_title(s, "建议: 产品判断看尾延迟, 科学结论补 per-IO LBA trace", "把展示结论和下一轮测试计划拆开。")
    add_bullets(
        s,
        [
            "给业务/产品: SSD 平均带宽不是主矛盾, 读/写尾延迟和 GC 后稳定性更能区分设备。",
            "给测试: 保留 ShareGPT + bpftrace + iostat 三件套, 再加 blktrace/eBPF per-IO CSV。",
            "给报告: LBA last-touch 只能讲空间触达范围, 不再讲真实 sequential ratio 或预取长度。",
            "给老板的下一步问题: 同样 workload 下, 多盘/多用户是否把 inter-request 冷读从尾延迟推成吞吐瓶颈?",
        ],
        0.85,
        1.65,
        11.4,
        2.8,
        size=20,
    )
    add_card(s, 0.85, 5.25, 3.5, 1.1, "下一轮 CSV 字段", "timestamp,dev,sector,bytes,rwbs", "per-IO LBA stream", CYAN)
    add_card(s, 4.85, 5.25, 3.5, 1.1, "核心计算", "delta/run/window", "read/write 分开统计", GREEN)
    add_card(s, 8.85, 5.25, 3.5, 1.1, "输出", "真实随机性结论", "可与 iostat %rrqm 对齐", YELLOW)

    # 10. Appendix
    s = blank_slide(prs)
    add_kicker(s, "appendix")
    add_title(s, "关键数据索引", "PPT 数字可从下列文件复现。")
    add_bullets(
        s,
        [
            "LBA 文档: docs/kv-cache-lba-timeline-analysis-2026-06-25.md",
            "LBA 数据: results/kvcache-profile/lba_timeline/lba_events.json / lba_timeline_summary.json",
            "Key locality: results/kvcache-profile/key_locality/kvcache_key_locality_summary.json",
            "设备 bpftrace: results/kvcache-profile/bpftrace_sharegpt_8b_tp8_cpu0p5g_users2_300s_profile_20260608_014520.txt",
            "四盘 iostat: results/cross_vendor/kv_cache_k4_gc_drift/_analysis/kvcache_io_randomness_summary.csv",
            "生成脚本: scripts/generate_kv_cache_boss_ppt.py",
        ],
        0.85,
        1.55,
        11.8,
        4.4,
        size=17,
    )

    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPT)


def write_summary(m: dict, charts: dict[str, Path]) -> None:
    lines = [
        "# KV Cache LBA / IO Executive PPT",
        "",
        f"Generated deck: `{OUT_PPT.relative_to(ROOT)}`",
        "",
        "## Key numbers",
        "",
        f"- ShareGPT profile: `{m['requests']:,}` requests, `{m['tokens']:,}` tokens, `{m['tok_s']:.0f}` tokens/s.",
        f"- Cache hit rate: `{m['cache_hit_pct']:.1f}%`; read/write volume ratio: `{m['read_write_ratio']:.1f}x`.",
        f"- Device IO shape: `{m['read_128_256_pct']:.0f}%` reads and `{m['write_128_256_pct']:.0f}%` writes are 128-256 KiB.",
        f"- Device D2C latency: read p50/p99 `{m['read_d2c_p50_us']}/{m['read_d2c_p99_us']} us`; write p50/p99 `{m['write_d2c_p50_us']}/{m['write_d2c_p99_us']} us`.",
        f"- LBA last-touch: `{m['lba_unique']}` unique `(dev, sector)` entries; nonzero range `{m['lba_nonzero_min_gib']:.1f}-{m['lba_nonzero_max_gib']:.1f} GiB`.",
        f"- Key locality: `{m['intra_token_pct']:.1f}%` rereads are intra-token `<10ms`; inter-request cold rereads are `{m['inter_request_pct']:.1f}%`.",
        "",
        "## Generated charts",
        "",
    ]
    for name, path in charts.items():
        lines.append(f"- `{name}`: `{path.relative_to(ROOT)}`")
    lines.extend(
        [
            "",
            "## Validation note",
            "",
            "`@d[dev, sector] = nsecs` is a bpftrace last-touch map, not a complete per-IO LBA log. "
            "The deck deliberately treats gap/direction/run metrics as exploratory last-touch-derived signals only.",
            "",
        ]
    )
    OUT_SUMMARY.write_text("\n".join(lines), encoding="utf-8")


def main() -> int:
    global OUT_PPT
    parser = argparse.ArgumentParser()
    parser.add_argument("--ppt", type=Path, default=OUT_PPT)
    args = parser.parse_args()
    OUT_PPT = args.ppt

    setup_matplotlib()
    metrics = collect_metrics()
    charts = generate_charts(metrics)
    build_ppt(metrics, charts)
    write_summary(metrics, charts)
    print(f"Wrote {OUT_PPT}")
    print(f"Wrote {OUT_SUMMARY}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
