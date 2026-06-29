#!/usr/bin/env python3
"""Generate a concise boss-facing AI SSD PPT deck.

The deck is intentionally one-message-per-slide. It reuses existing charts
where possible and generates a few simple schematic charts for product-design
slides.
"""

from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "presentations"
ASSET_DIR = ROOT / "docs" / "assets" / "ai-ssd-boss-ppt-20260630"
OUT_PPT = OUT_DIR / "ai-ssd-boss-deck-2026-06-30.pptx"
OUT_MD = OUT_DIR / "ai-ssd-boss-deck-2026-06-30.md"


BG = RGBColor(8, 12, 24)
PANEL = RGBColor(17, 25, 43)
PANEL_2 = RGBColor(24, 35, 58)
TEXT = RGBColor(238, 244, 255)
MUTED = RGBColor(155, 170, 195)
CYAN = RGBColor(41, 211, 255)
GREEN = RGBColor(63, 220, 146)
YELLOW = RGBColor(255, 201, 77)
RED = RGBColor(255, 91, 109)
PURPLE = RGBColor(176, 121, 255)
WHITE = RGBColor(255, 255, 255)


def setup_matplotlib() -> None:
    cjk_font_path = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
    if Path(cjk_font_path).exists():
        from matplotlib import font_manager as fm

        font_prop = fm.FontProperties(fname=cjk_font_path)
        fm.fontManager.addfont(cjk_font_path)
        plt.rcParams["font.family"] = "sans-serif"
        plt.rcParams["font.sans-serif"] = [font_prop.get_name(), "Noto Sans CJK SC", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    plt.rcParams["figure.facecolor"] = "#080c18"
    plt.rcParams["axes.facecolor"] = "#11192b"
    plt.rcParams["savefig.facecolor"] = "#080c18"
    plt.rcParams["text.color"] = "#eff6ff"
    plt.rcParams["axes.labelcolor"] = "#dce8ff"
    plt.rcParams["xtick.color"] = "#9fb0cc"
    plt.rcParams["ytick.color"] = "#9fb0cc"
    plt.rcParams["grid.color"] = "#33415f"


def savefig(path: Path) -> None:
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def chart_slc_context_buffer() -> Path:
    path = ASSET_DIR / "01_slc_context_buffer.png"
    fig, axes = plt.subplots(1, 2, figsize=(13.4, 5.9), gridspec_kw={"width_ratios": [1, 1.1]})

    ax = axes[0]
    ax.set_title("消费级动态 SLC: 峰值好看, 行为不可控", fontsize=15, fontweight="bold")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.axis("off")
    ax.add_patch(plt.Rectangle((8, 58), 35, 18, color="#29d3ff", alpha=0.9))
    ax.add_patch(plt.Rectangle((43, 58), 42, 18, color="#b079ff", alpha=0.65))
    ax.add_patch(plt.Rectangle((8, 28), 77, 18, color="#ff5b6d", alpha=0.72))
    ax.text(25.5, 67, "SLC burst", ha="center", va="center", fontsize=13, fontweight="bold")
    ax.text(64, 67, "TLC / QLC pool", ha="center", va="center", fontsize=13, fontweight="bold")
    ax.text(46.5, 37, "GC / fold 不可见, 可能污染前台 read", ha="center", va="center", fontsize=12)
    ax.annotate("短写很快", xy=(25, 78), xytext=(15, 91), arrowprops={"arrowstyle": "->", "color": "#29d3ff"}, color="#29d3ff", fontsize=12)
    ax.annotate("长稳态 cliff", xy=(47, 48), xytext=(46, 11), arrowprops={"arrowstyle": "->", "color": "#ff5b6d"}, color="#ff5b6d", fontsize=12)

    ax2 = axes[1]
    ax2.set_title("AI SSD pSLC context buffer: 可配置 + 可观测 + 可隔离", fontsize=15, fontweight="bold")
    ax2.set_xlim(0, 100)
    ax2.set_ylim(0, 100)
    ax2.axis("off")
    blocks = [
        (8, 68, 26, 15, "#3fdc92", "固定 pSLC\n64-256GiB"),
        (36, 68, 26, 15, "#29d3ff", "Hot KV / WAL\n低 tail"),
        (64, 68, 26, 15, "#b079ff", "TLC hot\ncontext"),
        (8, 40, 82, 15, "#ffc94d", "read-priority GC / fold 可暂停 / fold backlog telemetry"),
        (8, 15, 82, 15, "#4b5f86", "QLC / TLC cold context: RAG corpus, long memory, archive"),
    ]
    for x, y, w, h, color, label in blocks:
        ax2.add_patch(plt.Rectangle((x, y), w, h, color=color, alpha=0.88))
        ax2.text(x + w / 2, y + h / 2, label, ha="center", va="center", fontsize=11, fontweight="bold")
    ax2.annotate("host hint: KV / checkpoint / RAG / WAL", xy=(49, 84), xytext=(28, 96), arrowprops={"arrowstyle": "->", "color": "#eff6ff"}, fontsize=12)

    savefig(path)
    return path


def chart_product_matrix() -> Path:
    path = ASSET_DIR / "02_product_matrix.png"
    fig, ax = plt.subplots(figsize=(13.4, 5.9))
    ax.axis("off")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.text(2, 94, "AI SSD 产品分层: hot context 和 cold context 不能混为一谈", fontsize=18, fontweight="bold")

    rows = [
        ("TLC Performance", "Hot KV / decode miss / small checkpoint", "128KiB read P99, read-priority GC, GDS", "#3fdc92"),
        ("QLC Capacity", "RAG corpus / long memory / cold KV archive", "TB/$, cold read consistency, metadata pSLC", "#29d3ff"),
        ("Hybrid / System Kit", "AI PC / workstation / small AI server", "configurable pSLC, telemetry, placement hint", "#b079ff"),
    ]
    y = 72
    for name, target, spec, color in rows:
        ax.add_patch(plt.Rectangle((4, y - 10), 92, 18, color="#11192b", ec=color, lw=2))
        ax.text(8, y, name, fontsize=16, fontweight="bold", color=color, va="center")
        ax.text(35, y + 4, target, fontsize=12.5, color="#eff6ff", va="center")
        ax.text(35, y - 4, spec, fontsize=12, color="#9fb0cc", va="center")
        y -= 24
    ax.text(4, 8, "设计原则: TLC 做低 tail 热层, QLC 做容量冷层; pSLC 保护 metadata/WAL/KV hot write, 不承诺 QLC hot KV low tail。", fontsize=13, color="#ffc94d")
    savefig(path)
    return path


def chart_validation_roadmap() -> Path:
    path = ASSET_DIR / "03_validation_roadmap.png"
    fig, ax = plt.subplots(figsize=(13.4, 5.9))
    ax.axis("off")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.text(3, 93, "下一阶段: 先把测试标准做硬, 再谈产品规格", fontsize=18, fontweight="bold")

    stages = [
        ("P0", "Benchmark SOP", "block trace + activation gate + preconditioning", "#3fdc92"),
        ("P0", "Workload", "ShareGPT realistic + BurstGPT stress + 128KiB fio", "#29d3ff"),
        ("P1", "System", "Mooncake / LMCache clean offload, GDS vs non-GDS", "#b079ff"),
        ("P1", "Product", "TLC performance tier + QLC capacity tier", "#ffc94d"),
    ]
    x0 = 6
    for i, (level, title, body, color) in enumerate(stages):
        x = x0 + i * 23
        ax.add_patch(plt.Circle((x, 57), 7.5, color=color, alpha=0.95))
        ax.text(x, 57, level, ha="center", va="center", fontsize=15, fontweight="bold", color="#080c18")
        ax.text(x, 39, title, ha="center", fontsize=14, fontweight="bold", color=color)
        ax.text(x, 30, body, ha="center", fontsize=10.5, color="#dce8ff", wrap=True)
        if i < len(stages) - 1:
            ax.annotate("", xy=(x + 15, 57), xytext=(x + 8.5, 57), arrowprops={"arrowstyle": "->", "lw": 2, "color": "#9fb0cc"})
    ax.text(7, 10, "门禁: 没有 per-I/O LBA / SSD path proof / 3-run median 的图, 不能进入老板版性能结论。", fontsize=13, color="#ff5b6d")
    savefig(path)
    return path


def chart_gds_path() -> Path:
    path = ASSET_DIR / "04_gds_path.png"
    fig, axes = plt.subplots(1, 2, figsize=(13.4, 5.9))
    for ax in axes:
        ax.axis("off")
        ax.set_xlim(0, 100)
        ax.set_ylim(0, 100)
    axes[0].set_title("Non-GDS: SSD -> CPU DRAM -> GPU HBM", fontsize=15, fontweight="bold")
    axes[1].set_title("GDS: SSD -> GPU HBM, 但必须验证 direct path", fontsize=15, fontweight="bold")

    def box(ax, x, y, w, h, text, color):
        ax.add_patch(plt.Rectangle((x, y), w, h, color=color, alpha=0.9))
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=12, fontweight="bold")

    box(axes[0], 8, 58, 22, 16, "NVMe SSD", "#29d3ff")
    box(axes[0], 39, 58, 22, 16, "CPU DRAM\nbounce", "#ff5b6d")
    box(axes[0], 70, 58, 22, 16, "GPU HBM", "#3fdc92")
    axes[0].annotate("", xy=(38, 66), xytext=(31, 66), arrowprops={"arrowstyle": "->", "lw": 2, "color": "#eff6ff"})
    axes[0].annotate("", xy=(69, 66), xytext=(62, 66), arrowprops={"arrowstyle": "->", "lw": 2, "color": "#eff6ff"})
    axes[0].text(15, 30, "风险: CPU copy / NUMA / page cache / jitter 混入 SSD 结论", fontsize=12, color="#ffc94d")

    box(axes[1], 12, 58, 24, 16, "NVMe SSD", "#29d3ff")
    box(axes[1], 64, 58, 24, 16, "GPU HBM", "#3fdc92")
    axes[1].annotate("", xy=(63, 66), xytext=(37, 66), arrowprops={"arrowstyle": "->", "lw": 3, "color": "#3fdc92"})
    axes[1].text(18, 30, "门禁: gdscheck + cuFile logs + CPU util + fallback detection", fontsize=12, color="#ffc94d")
    axes[1].text(18, 20, "定位: P1/P2 预研, 不能先承诺必然提升 TTFT", fontsize=12, color="#ff5b6d")

    savefig(path)
    return path


def chart_hero_map() -> Path:
    path = ASSET_DIR / "00_hero_ai_ssd_map.png"
    fig, ax = plt.subplots(figsize=(13.4, 5.9))
    ax.axis("off")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)

    ax.text(4, 92, "AI SSD = LLM Context Memory Tier", fontsize=24, fontweight="bold")
    ax.text(4, 84, "从峰值带宽竞争, 转向长稳态 tail / GC / pSLC / GDS / telemetry", fontsize=13, color="#9fb0cc")

    center = (50, 50)
    ax.add_patch(plt.Circle(center, 15, color="#11192b", ec="#29d3ff", lw=2.4))
    ax.add_patch(plt.Circle(center, 18, fill=False, ec="#29d3ff", lw=1.2, alpha=0.35))
    ax.text(50, 54, "AI SSD", ha="center", va="center", fontsize=23, fontweight="bold", color="#eff6ff")
    ax.text(50, 45, "Context Buffer", ha="center", va="center", fontsize=12, color="#29d3ff", fontweight="bold")

    nodes = [
        ((17, 68), "KV Cache", "128KiB 随机读", "#ff5b6d"),
        ((23, 28), "pSLC", "可控上下文缓存", "#3fdc92"),
        ((50, 18), "GC QoS", "read-priority", "#ffc94d"),
        ((77, 28), "Telemetry", "可观测/可调度", "#b079ff"),
        ((83, 68), "GDS", "GPU 直连路径", "#29d3ff"),
    ]
    for (x, y), title, sub, color in nodes:
        ax.add_patch(plt.Circle((x, y), 9.2, color="#11192b", ec=color, lw=2))
        ax.text(x, y + 2.3, title, ha="center", va="center", fontsize=13, color=color, fontweight="bold")
        ax.text(x, y - 3.5, sub, ha="center", va="center", fontsize=9.5, color="#dce8ff")
        ax.annotate("", xy=(50 + (x - 50) * 0.62, 50 + (y - 50) * 0.62), xytext=(x, y), arrowprops={"arrowstyle": "->", "lw": 1.7, "color": color, "alpha": 0.88})

    # Signal strip
    signals = [
        ("Read jump", "95.1%", "#ff5b6d"),
        ("BurstGPT", "35K IOPS", "#29d3ff"),
        ("30min drift", "-44%", "#ffc94d"),
        ("SLC cache", "71-95GiB", "#3fdc92"),
        ("Mooncake", "41GiB offload", "#b079ff"),
    ]
    x = 5
    for label, value, color in signals:
        ax.add_patch(plt.Rectangle((x, 5), 17, 8.5, color="#11192b", ec=color, lw=1.2))
        ax.text(x + 1, 10.2, value, fontsize=13, fontweight="bold", color=color)
        ax.text(x + 1, 6.8, label, fontsize=9, color="#9fb0cc")
        x += 18.5

    savefig(path)
    return path


def chart_requirement_heatmap() -> Path:
    path = ASSET_DIR / "05_requirement_heatmap.png"
    rows = ["128KiB读P99", "长稳态GC", "Mixed R/W隔离", "pSLC上下文缓存", "GDS路径", "Telemetry"]
    cols = ["证据强度", "用户价值", "产品差异", "实现难度", "优先级"]
    values = [
        [5, 5, 4, 3, 5],
        [5, 5, 5, 4, 5],
        [4, 5, 5, 4, 5],
        [4, 4, 5, 4, 4],
        [3, 4, 4, 5, 3],
        [4, 4, 4, 3, 4],
    ]
    fig, ax = plt.subplots(figsize=(13.4, 5.9))
    im = ax.imshow(values, cmap="viridis", vmin=1, vmax=5)
    ax.set_xticks(range(len(cols)))
    ax.set_yticks(range(len(rows)))
    ax.set_xticklabels(cols, fontsize=12)
    ax.set_yticklabels(rows, fontsize=12)
    ax.set_title("AI SSD v1.0 需求热力图: P0 先做可验证指标", fontsize=18, fontweight="bold", pad=18)
    for i in range(len(rows)):
        for j in range(len(cols)):
            ax.text(j, i, str(values[i][j]), ha="center", va="center", fontsize=14, fontweight="bold", color="white")
    cbar = fig.colorbar(im, ax=ax, fraction=0.025, pad=0.02)
    cbar.set_label("1=低, 5=高", color="#dce8ff")
    ax.text(-0.45, 6.05, "解读: 128KiB读P99、长稳态GC、Mixed R/W隔离是P0; GDS是高端路线但需实测验证。", fontsize=12, color="#ffc94d")
    savefig(path)
    return path


def chart_evidence_to_product() -> Path:
    path = ASSET_DIR / "06_evidence_to_product.png"
    fig, ax = plt.subplots(figsize=(13.4, 5.9))
    ax.axis("off")
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 100)
    ax.text(3, 93, "从测试证据到产品定义: 不能跳过中间推理", fontsize=20, fontweight="bold")

    columns = [
        ("测试证据", ["真实LBA随机读", "30min token/s漂移", "SLC mixed R/W失效", "Mooncake path proof"], "#29d3ff"),
        ("设计需求", ["128KiB read tail", "read-priority GC", "可控pSLC", "activation gate"], "#3fdc92"),
        ("产品路线", ["TLC Hot Context", "QLC Cold Context", "GDS高端路径", "Benchmark Kit"], "#b079ff"),
    ]
    xs = [6, 38, 70]
    for x, (title, items, color) in zip(xs, columns):
        ax.add_patch(plt.Rectangle((x, 18), 24, 64, color="#11192b", ec=color, lw=2))
        ax.text(x + 12, 76, title, ha="center", fontsize=16, fontweight="bold", color=color)
        y = 62
        for item in items:
            ax.add_patch(plt.Rectangle((x + 3, y - 4.5), 18, 8, color="#18233a", ec="#33415f", lw=1))
            ax.text(x + 12, y, item, ha="center", va="center", fontsize=10.5, color="#eff6ff")
            y -= 12
    for x in [31, 63]:
        ax.annotate("", xy=(x + 5, 50), xytext=(x, 50), arrowprops={"arrowstyle": "->", "lw": 2.4, "color": "#ffc94d"})
    ax.text(6, 8, "汇报原则: 数据事实 -> 工程判断 -> 产品设计; 不用单次曲线直接承诺SKU或SLO。", fontsize=13, color="#ffc94d")
    savefig(path)
    return path


def make_charts() -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    return {
        "hero": chart_hero_map(),
        "slc": chart_slc_context_buffer(),
        "product": chart_product_matrix(),
        "roadmap": chart_validation_roadmap(),
        "gds": chart_gds_path(),
        "heatmap": chart_requirement_heatmap(),
        "flow": chart_evidence_to_product(),
    }


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_text(slide, text, x, y, w, h, size=20, color=TEXT, bold=False, align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text or " "
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.45, 0.45, 12.4, 0.55, size=24, color=TEXT, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.48, 1.05, 11.9, 0.35, size=12, color=MUTED)


def add_kicker(slide, text):
    add_text(slide, text.upper(), 0.45, 0.18, 8, 0.25, size=9, color=CYAN, bold=True)


def add_bullets(slide, items, x, y, w, h, size=16, color=TEXT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return tb


def add_card(slide, x, y, w, h, title, value, subtitle, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = color
    shape.line.width = Pt(1.1)
    add_text(slide, title, x + 0.15, y + 0.12, w - 0.3, 0.24, size=9, color=color, bold=True)
    add_text(slide, value, x + 0.15, y + 0.42, w - 0.3, 0.42, size=20, color=TEXT, bold=True)
    add_text(slide, subtitle, x + 0.15, y + 0.88, w - 0.3, 0.36, size=8.5, color=MUTED)


def add_bottom_line(slide, text):
    add_text(slide, text, 0.55, 6.88, 12.2, 0.28, size=11, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)


def add_image(slide, path: Path, x, y, w, h=None):
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h) if h else None)


def image(path: str) -> Path:
    p = ROOT / path
    if not p.exists():
        raise FileNotFoundError(p)
    return p


def build_ppt(extra_charts: dict[str, Path]) -> list[dict[str, str]]:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    notes: list[dict[str, str]] = []

    def note(title: str, talk: str):
        notes.append({"title": title, "talk": talk})

    # 1
    s = blank_slide(prs)
    add_text(s, "AI SSD 产品预研", 0.65, 0.72, 8.8, 0.72, size=34, bold=True)
    add_text(s, "从 KV Cache 真实 I/O 到 pSLC Context Buffer 的产品设计", 0.7, 1.42, 11.5, 0.4, size=16, color=MUTED)
    add_image(s, extra_charts["hero"], 0.55, 1.9, 12.25, 4.72)
    add_bottom_line(s, "一句话: 把 SSD 从容量盘升级为 LLM serving 的 context memory tier。")
    note("AI SSD 产品预研", "开场只讲一个判断：AI SSD 的卖点不是峰值顺序带宽，而是上下文数据在长时间服务中的稳定性。")

    # 2
    s = blank_slide(prs)
    add_kicker(s, "why now")
    add_title(s, "为什么需要 AI SSD: KV Cache 正在变成模型的上下文记忆层", "长上下文、多轮会话、Agent memory 会把 HBM/DRAM 放不下的上下文下沉到 SSD。")
    add_image(s, image("docs/mooncake-source-research/charts/01_four_tier_cache_hierarchy.png"), 0.65, 1.55, 6.4, 4.55)
    add_bullets(
        s,
        [
            "HBM 最快但最贵, 不能无限放上下文。",
            "DRAM 可以缓冲, 但长会话和多用户仍会溢出。",
            "SSD 成为 context memory tier: 它必须服务 TTFT 和 read tail, 不是只存文件。",
        ],
        7.35,
        2.0,
        5.2,
        2.6,
        size=17,
    )
    add_bottom_line(s, "产品含义: AI SSD 要用系统指标定义价值, 例如 TTFT、cache hit、read P99、GC cliff。")
    note("为什么需要 AI SSD", "这页解释战略背景：KV cache 从临时显存状态变成可迁移、可复用的 context memory。")

    # 3
    s = blank_slide(prs)
    add_kicker(s, "real io")
    add_title(s, "真实 I/O 画像: decode read 随机, write 多数连续", "真实 block trace 显示, KV cache 不能再笼统说成“随机大块 IO”。")
    add_image(s, image("docs/assets/kv-cache-real-io/04_delta_signature.png"), 0.55, 1.42, 7.1, 4.95)
    add_card(s, 7.9, 1.62, 2.25, 1.05, "Read jump", "95.1%", ">=100 MiB", RED)
    add_card(s, 10.35, 1.62, 2.25, 1.05, "Write contig", "75.1%", "exact contiguous", GREEN)
    add_card(s, 7.9, 3.0, 2.25, 1.05, "Request", "128KiB", "dominant size", CYAN)
    add_card(s, 10.35, 3.0, 2.25, 1.05, "LBA span", "389GiB", "wide touch", PURPLE)
    add_bullets(s, ["SSD 选型主压力是 decode random read tail。", "写入不是无关紧要: 它通过 GC / fold 污染前台 read。"], 8.0, 4.55, 4.6, 1.1, size=15)
    add_bottom_line(s, "产品含义: 重点优化 128KiB random read P99/P999, 不是只报 4K random 或顺序峰值。")
    note("真实 I/O 画像", "数据页：读写分裂。read 是随机压力源，write 是 GC 污染源。")

    # 4
    s = blank_slide(prs)
    add_kicker(s, "workload taxonomy")
    add_title(s, "三类 workload 不能互相替代", "Synthetic、ShareGPT、BurstGPT 分别回答不同问题。")
    add_image(s, image("docs/assets/io-three-way-comparison/01_signal_dashboard.png"), 0.55, 1.42, 7.6, 4.95)
    add_bullets(
        s,
        [
            "Synthetic fio: 设备能力标定, 可扫 QD/P99, 不能证明真实 LBA。",
            "ShareGPT: 真实聊天 replay, 压力较轻, 适合体验验证。",
            "BurstGPT: 更重更随机, 适合 SSD stress baseline。",
        ],
        8.45,
        1.8,
        4.25,
        2.2,
        size=15.5,
    )
    add_card(s, 8.45, 4.45, 1.95, 1.0, "ShareGPT", "14K", "block IOPS", CYAN)
    add_card(s, 10.75, 4.45, 1.95, 1.0, "BurstGPT", "35K", "block IOPS", RED)
    add_bottom_line(s, "产品含义: 认证标准至少要包含 realistic replay + stress replay + fio baseline。")
    note("三类 workload", "这页避免老板把 fio、ShareGPT、BurstGPT混为一谈。")

    # 5
    s = blank_slide(prs)
    add_kicker(s, "long steady")
    add_title(s, "短时峰值不是 AI SSD 能力: 30 分钟 token/s 下降 44%", "长稳态下 SSD GC / 写放大 / tail latency 会传导到 serving 吞吐。")
    add_image(s, image("docs/assets/charts/k5_70b_6users_30min_token_rate.png"), 0.62, 1.45, 7.65, 4.85)
    add_card(s, 8.55, 1.72, 2.0, 1.0, "Peak", "3113", "token/s", GREEN)
    add_card(s, 10.8, 1.72, 2.0, 1.0, "Tail", "1752", "token/s", RED)
    add_card(s, 8.55, 3.08, 4.25, 1.0, "Drift", "-44%", "30min 内吞吐下滑", YELLOW)
    add_bullets(s, ["5 分钟短测会高估能力。", "老板应看 30/60/120min token/s 和 P99 drift。"], 8.65, 4.55, 4.1, 0.95, size=15)
    add_bottom_line(s, "产品含义: AI SSD 需要 predictable GC, 否则用户看到的是 TTFT 和吞吐抖动。")
    note("长稳态 token/s", "这页按你之前问的 token/s 曲线做：一页只讲长稳态会拖垮吞吐。")

    # 6
    s = blank_slide(prs)
    add_kicker(s, "gc cliff")
    add_title(s, "GC cliff 会改变短测选型结论", "Biwin 短测领先, 20-30min 后与 Seagate 收敛; 这说明 firmware/GC 比峰值更重要。")
    add_image(s, image("docs/assets/charts/07_long_drift_compare.png"), 0.58, 1.48, 7.65, 4.85)
    add_bullets(
        s,
        [
            "Biwin: 3.14GB/s 短测 -> 1.57GB/s 30min。",
            "Seagate: 2.34GB/s 短测 -> 1.54GB/s 30min。",
            "两者长稳态仅差约 2%; 短时冠军不等于服务冠军。",
        ],
        8.55,
        1.85,
        4.05,
        2.35,
        size=15,
    )
    add_card(s, 8.55, 4.65, 1.95, 1.0, "Biwin", "2.9min", "GC cliff", RED)
    add_card(s, 10.8, 4.65, 1.95, 1.0, "Seagate", "8.1min", "GC cliff", GREEN)
    add_bottom_line(s, "产品含义: Spec 里必须加入 GC cliff time、drop、recovery 和 late-window P99。")
    note("GC cliff", "这页讲选型逻辑：短测排名会失真。")

    # 7
    s = blank_slide(prs)
    add_kicker(s, "slc cache")
    add_title(s, "SLC cache 不是答案本身: 要改造成 pSLC context buffer", "消费级 SLC 为短写 burst 设计; AI 需要可预测、可隔离、可观测。")
    add_image(s, extra_charts["slc"], 0.58, 1.5, 12.15, 4.95)
    add_bottom_line(s, "产品含义: 不是“更大 SLC cache”, 而是固定/可配置 pSLC + read-priority GC + telemetry。")
    note("SLC cache", "这页回应 TLC/QLC 消费级 SLC cache：方向有价值，但要从 burst buffer 变成 context buffer。")

    # 8
    s = blank_slide(prs)
    add_kicker(s, "mooncake proof")
    add_title(s, "系统级 offload 必须先证明 SSD path 真的触发", "配置名叫 SSD 没用; 要有 root、enable、read store、O_DIRECT、文件增长证据。")
    add_image(s, image("docs/assets/mooncake-ssd-offload-final-formal-20260629/01_overall_performance_local.png"), 0.58, 1.45, 7.45, 4.85)
    add_card(s, 8.3, 1.62, 2.05, 1.0, "TTFT", "-17.2%", "+SSD vs Mooncake", GREEN)
    add_card(s, 10.65, 1.62, 2.05, 1.0, "Cache hit", "67.8%", "+Mooncake+SSD", CYAN)
    add_card(s, 8.3, 2.95, 2.05, 1.0, "Offload", "41GiB", "files on SSD", PURPLE)
    add_card(s, 10.65, 2.95, 2.05, 1.0, "O_DIRECT", "1341", "events", YELLOW)
    add_bullets(s, ["当前结论: path 已触发, 趋势成立。", "边界: 仍有 insufficient space, 不是 clean production benchmark。"], 8.35, 4.55, 4.25, 1.0, size=14.5)
    add_bottom_line(s, "产品含义: AI SSD benchmark 必须有 activation gate, 否则性能图没有 SSD 归因价值。")
    note("Mooncake path proof", "这页讲方法论：系统级收益必须先证明数据真的走 SSD。")

    # 9
    s = blank_slide(prs)
    add_kicker(s, "gds path")
    add_title(s, "GDS 是高端方向, 但不能先承诺收益", "Non-GDS 会引入 CPU bounce buffer; GDS 需要证明 direct path 没 fallback。")
    add_image(s, extra_charts["gds"], 0.65, 1.55, 12.0, 4.8)
    add_bottom_line(s, "产品含义: GDS readiness 应作为 P1/P2, 指标包括 CPU util、TTFT、read tail、fallback detection。")
    note("GDS", "这页讲未来路径：值得做，但必须实测，不能把配置当结果。")

    # 10
    s = blank_slide(prs)
    add_kicker(s, "product design")
    add_title(s, "产品设计: TLC 热层 + QLC 冷层, pSLC 做 context buffer", "把介质、缓存和 workload 分层, 不要用一个 SSD 口号覆盖所有 AI 场景。")
    add_image(s, extra_charts["product"], 0.62, 1.5, 12.05, 4.85)
    add_bottom_line(s, "产品含义: Hot KV 需要 TLC low tail; QLC 更适合 RAG/cold memory; pSLC 保护 metadata/WAL/hot write。")
    note("产品设计", "这页把设计收敛成三类产品路线。")

    # 11
    s = blank_slide(prs)
    add_kicker(s, "requirements")
    add_title(s, "AI SSD v1.0 需求: 先定义可验证的指标", "把普通 SSD spec 改成 AI workload scorecard。")
    add_image(s, extra_charts["heatmap"], 0.68, 1.52, 12.0, 4.95)
    add_bottom_line(s, "产品含义: 没有 tail、GC、path proof 的 AI SSD 指标不可信。")
    note("AI SSD v1.0 需求", "这页给出需求清单，适合老板拍下一阶段资源。")

    # 12
    s = blank_slide(prs)
    add_kicker(s, "next step")
    add_title(s, "下一步: 用标准化实验把预研变成产品判断", "先证明指标稳定, 再进入样机和规格定义。")
    add_image(s, extra_charts["flow"], 0.65, 1.52, 12.0, 4.75)
    add_bottom_line(s, "管理层决策: 继续投入 AI SSD 预研, 但阶段目标是测试标准 + 样机方向, 不是立刻承诺最终 SKU。")
    note("下一步", "收尾：建议继续投入，但不要过早承诺型号或最终 SLO。")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPT)
    return notes


def write_summary(notes: list[dict[str, str]]) -> None:
    lines = [
        "# AI SSD Boss Deck",
        "",
        f"Generated deck: `{OUT_PPT.relative_to(ROOT)}`",
        "",
        "## Slide Outline",
        "",
    ]
    for i, item in enumerate(notes, 1):
        lines.append(f"### {i}. {item['title']}")
        lines.append("")
        lines.append(item["talk"])
        lines.append("")
    OUT_MD.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    setup_matplotlib()
    charts = make_charts()
    notes = build_ppt(charts)
    write_summary(notes)
    print(f"Wrote {OUT_PPT}")
    print(f"Wrote {OUT_MD}")


if __name__ == "__main__":
    main()
